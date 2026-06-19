#!/usr/bin/env python3
"""Genera una presentacion (.pptx) que explica e compara las 4 tecnicas
modernas de embedding de subgrafos urbanos implementadas en el dashboard:

    1. GL2Vec    — Chen & Koga, ICONIP 2019
    2. FEATHER-G — Rozemberczki & Sarkar, CIKM 2020
    3. DHC-E     — Wang, Deng, Lu & Chen, 2022 (arXiv:2108.02113)
    4. GCN       — Kipf & Welling, 2017, via el ST-GCN de Fan, Hu & Hu, 2025

Cada una se proyecta a 2-D con UMAP y se expone como un espacio del panel UMAP
de la pagina (embedGL2Vec / embedFeather / embedDHCE / embedGCN). Estilo y
paleta reusan los de make_ppt.py para mantener la coherencia visual.

Uso:
    python make_ppt_embeddings.py [salida.pptx]
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ─── Paleta (igual que make_ppt.py) ──────────────────────────────────────────
INK      = RGBColor(0x1A, 0x1F, 0x2B)
ACCENT   = RGBColor(0xFF, 0x6B, 0x4A)   # naranja (GL2Vec)
ACCENT2  = RGBColor(0x4A, 0x90, 0xD9)   # azul (FEATHER)
MUTED    = RGBColor(0x6B, 0x73, 0x80)
LIGHT    = RGBColor(0xF4, 0xF5, 0xF7)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GREEN    = RGBColor(0x3F, 0xB6, 0x80)   # verde (DHC-E)
PURPLE   = RGBColor(0x8B, 0x5C, 0xF6)   # violeta (GCN)

# Color por tecnica (para mantener un codigo visual consistente).
TECH_COLORS = {"GL2Vec": ACCENT, "FEATHER-G": ACCENT2, "DHC-E": GREEN, "GCN": PURPLE}

W, H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def box(s, x, y, w, h):
    tb = s.shapes.add_textbox(x, y, w, h)
    tb.text_frame.word_wrap = True
    return tb


def rect(s, x, y, w, h, color, line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def setp(p, text, size, color, bold=False, align=PP_ALIGN.LEFT, font="Calibri", space=6):
    p.text = text or " "  # an empty string creates no run; keep one to style
    p.alignment = align
    p.space_after = Pt(space)
    r = p.runs[0]
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = font
    return p


def add_para(tf, text, size, color, bold=False, align=PP_ALIGN.LEFT, space=6, font="Calibri"):
    p = tf.add_paragraph()
    return setp(p, text, size, color, bold, align, font, space)


def bg(s, color):
    rect(s, 0, 0, W, H, color)


def bullets(tf, items, size=15, color=INK, gap=8):
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        lvl, txt = (it if isinstance(it, tuple) else (0, it))
        p.level = lvl
        setp(p, ("•  " if lvl == 0 else "–  ") + txt, size, color, space=gap)


def header(s, kicker, title, accent=ACCENT):
    rect(s, Inches(0), Inches(0), Inches(0.18), H, accent)
    t = box(s, Inches(0.6), Inches(0.35), Inches(12), Inches(0.4))
    setp(t.text_frame.paragraphs[0], kicker, 13, accent, bold=True)
    t2 = box(s, Inches(0.55), Inches(0.7), Inches(12.2), Inches(0.9))
    setp(t2.text_frame.paragraphs[0], title, 30, INK, bold=True)


def paper_box(s, x, y, w, cite, accent):
    """Caja de cita del paper (referencia)."""
    rect(s, x, y, w, Inches(0.92), INK)
    tb = box(s, x + Inches(0.2), y + Inches(0.1), w - Inches(0.4), Inches(0.75))
    setp(tb.text_frame.paragraphs[0], "PAPER", 11, accent, bold=True, space=2)
    add_para(tb.text_frame, cite, 12.5, WHITE, space=0)


# ─── 1. Portada ──────────────────────────────────────────────────────────────
s = slide()
bg(s, INK)
rect(s, Inches(0.9), Inches(2.3), Inches(2.4), Inches(0.09), ACCENT)
t = box(s, Inches(0.85), Inches(2.55), Inches(11.6), Inches(2))
setp(t.text_frame.paragraphs[0],
     "Embeddings de subgrafos urbanos", 40, WHITE, bold=True)
add_para(t.text_frame, "3 tecnicas recientes · proyectadas en el UMAP de la pagina",
         22, ACCENT, space=4)
t2 = box(s, Inches(0.9), Inches(4.7), Inches(11.5), Inches(2))
setp(t2.text_frame.paragraphs[0],
     "Como convertir cada subgrafo de la red vial en un vector y verlo en 2-D.",
     15, MUTED)
for tech, col in TECH_COLORS.items():
    add_para(t2.text_frame, "■  " + tech, 15, col, bold=True, space=3)
add_para(t2.text_frame,
         "preprocess_hotspots.py · embed_subgraphs.py · components/UmapPanel.tsx",
         13, ACCENT2, space=2)

# ─── 2. Contexto / problema ──────────────────────────────────────────────────
s = slide()
bg(s, WHITE)
header(s, "CONTEXTO", "Por que embeddings de subgrafos")
tf = box(s, Inches(0.7), Inches(1.9), Inches(6.0), Inches(5)).text_frame
bullets(tf, [
    "El pipeline genera, por mes, las top-20 cuencas (subgrafos) de concentracion criminal sobre la red vial de Chicago.",
    "Para responder '¿que zonas se parecen?' necesitamos describir cada subgrafo con un vector comparable.",
    "Un embedding de grafo completo (whole-graph) resume la forma de la sub-red vial en un vector de longitud fija.",
    "UMAP reduce ese vector a 2-D: subgrafos parecidos quedan cerca en el scatter.",
    "Clave de diseño: el crimen NO entra al embedding. Asi 'vecinos con criminalidad distinta' es un hallazgo, no algo impuesto.",
], size=15, gap=10)
rect(s, Inches(7.05), Inches(2.0), Inches(5.5), Inches(4.3), LIGHT)
d = box(s, Inches(7.3), Inches(2.2), Inches(5.1), Inches(4)).text_frame
setp(d.paragraphs[0], "De subgrafo a punto 2-D", 16, ACCENT, bold=True, space=10)
for ln in ["sub-red vial del basin", "↓", "descriptor de grafo (vector)",
           "↓  (4 tecnicas distintas)", "GL2Vec · FEATHER · DHC-E · GCN", "↓",
           "UMAP → [x, y] en [0,1]", "↓", "punto en el panel UMAP"]:
    is_arrow = "↓" in ln
    add_para(d, ln, 12 if is_arrow else 15, MUTED if is_arrow else INK,
             bold=("GL2Vec" in ln or "UMAP →" in ln),
             align=PP_ALIGN.CENTER, space=4)

# ─── 3. Marco comun ──────────────────────────────────────────────────────────
s = slide()
bg(s, WHITE)
header(s, "MARCO COMUN", "Mismo esqueleto, 3 descriptores distintos")
chain = [("sub-red\nvial", LIGHT, INK),
         ("descriptor\n(tecnica)", ACCENT, WHITE),
         ("z-score\n(estandarizar)", LIGHT, INK),
         ("UMAP\n(cosine)", ACCENT2, WHITE),
         ("[x, y]\nen [0,1]", GREEN, WHITE)]
x = Inches(0.55)
for i, (txt, col, fg) in enumerate(chain):
    rect(s, x, Inches(2.0), Inches(2.15), Inches(1.3), col)
    b = box(s, x + Inches(0.1), Inches(2.2), Inches(1.95), Inches(0.95))
    for j, ln in enumerate(txt.split("\n")):
        p = b.text_frame.paragraphs[0] if j == 0 else b.text_frame.add_paragraph()
        setp(p, ln, 13, fg, bold=True, align=PP_ALIGN.CENTER, space=2)
    if i < 4:
        ar = box(s, x + Inches(2.13), Inches(2.3), Inches(0.45), Inches(0.7))
        setp(ar.text_frame.paragraphs[0], "→", 26, MUTED, align=PP_ALIGN.CENTER)
    x += Inches(2.55)
tf = box(s, Inches(0.7), Inches(3.7), Inches(12), Inches(3)).text_frame
bullets(tf, [
    "Lo unico que cambia entre las 4 tecnicas es como se construye el descriptor del subgrafo.",
    "El resto es identico: estandarizacion por columna, UMAP (n_neighbors=15, min_dist=0.1, metric=cosine, random_state=42) y escalado min-max a [0,1].",
    "Las 4 dan un vector por subgrafo (whole-graph): invariantes a permutacion y comparables entre subgrafos de distinto tamaño.",
    "Implementadas self-contained (numpy / networkx / gensim), sin karateclub (incompatible con Python 3.13).",
    "Si umap-learn falta, hay fallback a PCA(2); si gensim falta, GL2Vec cae a una traza NetLSD.",
], size=14.5, gap=9)

# ─── 4-6. Una slide por tecnica ──────────────────────────────────────────────
def tech_slide(kicker, title, accent, cite, idea, impl, captures):
    s = slide()
    bg(s, WHITE)
    header(s, kicker, title, accent)
    paper_box(s, Inches(0.7), Inches(1.75), Inches(11.9), cite, accent)
    # Columna izquierda: idea
    li = box(s, Inches(0.7), Inches(2.95), Inches(5.7), Inches(0.4))
    setp(li.text_frame.paragraphs[0], "LA IDEA", 13, accent, bold=True)
    tf = box(s, Inches(0.7), Inches(3.35), Inches(5.7), Inches(3.6)).text_frame
    bullets(tf, idea, size=14, gap=8)
    # Columna derecha arriba: como lo implementamos
    ri = box(s, Inches(6.8), Inches(2.95), Inches(5.8), Inches(0.4))
    setp(ri.text_frame.paragraphs[0], "COMO LO IMPLEMENTAMOS AQUI", 13, accent, bold=True)
    tf2 = box(s, Inches(6.8), Inches(3.35), Inches(5.8), Inches(2.2)).text_frame
    bullets(tf2, impl, size=13.5, gap=6)
    # Columna derecha abajo: que captura
    rc = box(s, Inches(6.8), Inches(5.55), Inches(5.8), Inches(0.4))
    setp(rc.text_frame.paragraphs[0], "QUE CAPTURA / FORTALEZA", 13, accent, bold=True)
    tf3 = box(s, Inches(6.8), Inches(5.95), Inches(5.8), Inches(1.3)).text_frame
    bullets(tf3, captures, size=13.5, gap=5)


tech_slide(
    "TECNICA 1 · GL2Vec", "GL2Vec — line graph + features de aristas", ACCENT,
    "Chen, H. & Koga, H. (2019). GL2vec: Graph Embedding Enriched by Line Graphs "
    "with Edge Features. ICONIP 2019, LNCS 11953, pp. 3-14.",
    idea=[
        "Extiende graph2vec, que solo mira los NODOS del grafo.",
        "Embebe tambien el line graph L(G): cada arista (segmento vial) se vuelve un nodo; dos segmentos son adyacentes si comparten una interseccion.",
        "El vector del subgrafo = concatenacion de graph2vec(G) y graph2vec(L(G)).",
        "L(G) expone la estructura que vive en las aristas (como encadenan los segmentos), que graph2vec pierde.",
    ],
    impl=[
        "Subgrafo inducido G y su nx.line_graph(G).",
        "Documentos Weisfeiler-Lehman (WL-2) de G y de L(G), con vocabularios separados.",
        (1, "2 modelos Doc2Vec (PV-DBOW, 32-D c/u) → hstack = firma 64-D."),
        "Variante estructural (etiqueta = grado): la clase de via por-arista no se guarda por separado.",
    ],
    captures=[
        "Aprende motivos recurrentes de la sub-red (cruces en T, fondos de saco) a lo largo de todo el corpus de 480 subgrafos.",
        "Sensible a la organizacion de los segmentos, no solo de las intersecciones.",
    ],
)

tech_slide(
    "TECNICA 2 · FEATHER-G", "FEATHER-G — funciones caracteristicas sobre random walks", ACCENT2,
    "Rozemberczki, B. & Sarkar, R. (2020). Characteristic Functions on Graphs: "
    "Birds of a Feather, from Statistical Descriptors to Parametric Models. CIKM 2020.",
    idea=[
        "Describe la DISTRIBUCION de un feature de nodo (aqui el grado) a varias escalas.",
        "Usa la funcion caracteristica E[e^{iθx}] evaluada sobre las probabilidades de transicion de random walks de r pasos.",
        "Parte real e imaginaria por nodo → se agrupan (pooling por media) en un descriptor del grafo.",
        "Determinista (sin entrenamiento) y se prueba que da el mismo vector a grafos isomorfos.",
    ],
    impl=[
        "P = D⁻¹A (matriz de transicion). Escalas r = 1..5.",
        "5 puntos de evaluacion θ; por nodo: P^r·cos(θx) y P^r·sin(θx).",
        (1, "Pooling por media → firma de 2·5·5 = 50-D."),
        "Basins sin aristas: se usa la funcion caracteristica de x directamente.",
    ],
    captures=[
        "Resume como se difunde la estructura del vecindario a distintas distancias.",
        "Robusto a corrupcion de datos; rapido y sin hiperparametros que ajustar (salvo escalas/θ).",
    ],
)

tech_slide(
    "TECNICA 3 · DHC-E", "DHC-E — Grado → H-index → Coreness + entropia", GREEN,
    "Wang, H., Deng, Y., Lu, L. & Chen, G. (2022). Hyperparameter-free and "
    "Explainable Whole Graph Embedding. arXiv:2108.02113.",
    idea=[
        "Itera el operador DHC: el valor de cada nodo se reemplaza por el H-index de los valores de sus vecinos.",
        "Por el teorema DHC, la cadena Grado → H-index → … converge a la Coreness del nodo.",
        "En cada iteracion se mide la entropia de Shannon del histograma de valores.",
        "El embedding es la secuencia de entropias [E0, E1, …]: sin hiperparametros y explicable.",
    ],
    impl=[
        "Iteracion 0 = grado de cada nodo del subgrafo.",
        "Se itera el H-index hasta el punto fijo (coreness).",
        (1, "Entropia normalizada por iteracion → secuencia padded a 16-D."),
        "Es la tecnica mas barata: solo grados y H-index, sin algebra lineal.",
    ],
    captures=[
        "Resume la jerarquia de centralidad/coreness del subgrafo en pocas cifras interpretables.",
        "Cada dimension tiene significado (entropia en el nivel t de la cadena).",
    ],
)

tech_slide(
    "TECNICA 4 · GCN", "GCN — convolucion de grafo (modulo espacial ST-GCN)", PURPLE,
    "Modulo espacial de Fan, Y.; Hu, X.; Hu, J. (2025). Research on a Crime "
    "Spatiotemporal Prediction Method Integrating Informer and ST-GCN: A Case "
    "Study of Four Crime Types in Chicago. Big Data Cogn. Comput. 9(7), 179. "
    "GCN: Kipf & Welling, ICLR 2017.",
    idea=[
        "El paper predice crimen en Chicago combinando Informer (tiempo) y ST-GCN (espacio).",
        "Su modulo espacial embebe los nodos del grafo con una GCN (Kipf & Welling).",
        "Cada capa: H' = ReLU( D̃^(-1/2)(A+I)D̃^(-1/2) · H · W ): cada nodo mezcla los features de sus vecinos por la adyacencia normalizada.",
        "Apilar capas da a cada nodo un embedding de su vecindario a k saltos.",
    ],
    impl=[
        "Adyacencia normalizada simetrica con auto-bucles; 2 capas (receptivo de 2 saltos).",
        "Features de nodo solo estructurales: [1, grado, log(1+grado)].",
        (1, "Mean-pooling de los nodos → descriptor de 16-D por subgrafo."),
        "GCN SIN entrenar (pesos aleatorios fijos): el crimen no puede formar el layout, asi que la usamos como extractor estructural.",
    ],
    captures=[
        "Difusion del vecindario por la adyacencia normalizada: estructura local a k saltos.",
        "Es la unica de las 4 basada en GNN; conecta directamente con el paper de prediccion.",
    ],
)

# ─── 7. Comparacion ──────────────────────────────────────────────────────────
s = slide()
bg(s, WHITE)
header(s, "COMPARACION", "Las 3 tecnicas, lado a lado")

rows = [
    ["", "GL2Vec", "FEATHER-G", "DHC-E", "GCN"],
    ["Paradigma", "Aprendido (Doc2Vec)", "Estadistico (random walk)", "Jerarquico (grado/entropia)", "GNN (conv. de grafo)"],
    ["Entrada", "G + line graph L(G)", "Feature nodo + P=D⁻¹A", "Cadena Grado→H-index", "Â=D̃^-½(A+I)D̃^-½"],
    ["¿Aprende?", "Si (corpus completo)", "No (deterministico)", "No (deterministico)", "No (pesos fijos)"],
    ["Invariante iso.", "Si (WL)", "Si (demostrado)", "Si", "Si (equivariante)"],
    ["Dim. descriptor", "64-D", "50-D", "16-D", "16-D"],
    ["Coste", "Medio (Doc2Vec)", "Bajo-medio", "Muy bajo", "Bajo"],
    ["Fuerte en", "Motivos de segmentos", "Difusion multi-escala", "Jerarquia / explicab.", "Vecindario a k saltos"],
    ["Limite", "Necesita gensim", "Eleccion de escalas/θ", "Solo grado como señal", "Sin entrenar; solo grado"],
]
nrows, ncols = len(rows), 5
gtab = s.shapes.add_table(nrows, ncols, Inches(0.5), Inches(1.75),
                          Inches(12.35), Inches(4.9)).table
gtab.columns[0].width = Inches(1.75)
for c in range(1, 5):
    gtab.columns[c].width = Inches(2.65)

col_accent = [INK, ACCENT, ACCENT2, GREEN, PURPLE]
for r in range(nrows):
    for c in range(ncols):
        cell = gtab.cell(r, c)
        cell.margin_left = Inches(0.07)
        cell.margin_right = Inches(0.05)
        cell.margin_top = Inches(0.03)
        cell.margin_bottom = Inches(0.03)
        tfc = cell.text_frame
        tfc.word_wrap = True
        p = tfc.paragraphs[0]
        txt = rows[r][c]
        is_head = (r == 0)
        is_label = (c == 0)
        if is_head:
            cell.fill.solid(); cell.fill.fore_color.rgb = col_accent[c]
            setp(p, txt, 13, WHITE, bold=True, align=PP_ALIGN.CENTER, space=0)
        elif is_label:
            cell.fill.solid(); cell.fill.fore_color.rgb = INK
            setp(p, txt, 11, WHITE, bold=True, space=0)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else LIGHT
            setp(p, txt, 10.5, INK, space=0)

# ─── 7b. HDBSCAN sobre cada UMAP ─────────────────────────────────────────────
s = slide()
bg(s, WHITE)
header(s, "CLUSTERING", "HDBSCAN sobre cada UMAP", PURPLE)
tf = box(s, Inches(0.7), Inches(1.9), Inches(6.1), Inches(5)).text_frame
bullets(tf, [
    "Para CADA tecnica corremos HDBSCAN sobre su propio layout 2-D del UMAP.",
    "HDBSCAN agrupa por densidad: no fija el numero de clusters y marca como 'ruido' los puntos que no caen en ninguna region densa.",
    "Clusterizar sobre el mismo 2-D que se ve garantiza que los colores pintados coincidan con los grupos visibles en el scatter.",
    "Parametros: min_cluster_size = 8, min_samples = 4 (sklearn.cluster.HDBSCAN).",
    "En la pagina: boton de color 'Clusteres' -> cada punto se pinta por su cluster; el ruido queda en gris.",
    "Cada tecnica produce su propia particion: comparar las 4 muestra que 'grupos de subgrafos parecidos' depende del descriptor.",
], size=14.5, gap=9)

# Recuento de clusters por tecnica (resultado de la corrida actual).
rect(s, Inches(7.05), Inches(2.0), Inches(5.5), Inches(4.3), LIGHT)
hd = box(s, Inches(7.3), Inches(2.15), Inches(5.1), Inches(0.5))
setp(hd.text_frame.paragraphs[0], "Clusteres encontrados (480 subgrafos)", 15, PURPLE, bold=True)
counts = [("GL2Vec", ACCENT, "8 clusteres · 0 ruido"),
          ("FEATHER-G", ACCENT2, "14 clusteres · 25 ruido"),
          ("DHC-E", GREEN, "10 clusteres · 8 ruido"),
          ("GCN", PURPLE, "21 clusteres · 18 ruido")]
y = Inches(2.85)
for tech, col, txt in counts:
    rect(s, Inches(7.3), y, Inches(0.32), Inches(0.32), col)
    tb = box(s, Inches(7.75), y - Inches(0.04), Inches(4.6), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    setp(p, tech + ":  ", 14, INK, bold=True, space=0)
    r = p.add_run(); r.text = txt
    r.font.size = Pt(14); r.font.color.rgb = MUTED; r.font.name = "Calibri"
    y += Inches(0.62)
note = box(s, Inches(7.3), Inches(5.5), Inches(5.0), Inches(0.8))
setp(note.text_frame.paragraphs[0],
     "Mas clusteres/ruido = la tecnica separa los subgrafos en grupos mas finos.",
     12.5, MUTED)

# ─── 8. Cierre / como se ven en la pagina ────────────────────────────────────
s = slide()
bg(s, INK)
rect(s, Inches(0.9), Inches(1.3), Inches(2.4), Inches(0.09), ACCENT)
t = box(s, Inches(0.85), Inches(1.55), Inches(11.5), Inches(1))
setp(t.text_frame.paragraphs[0], "En la pagina", 34, WHITE, bold=True)
tf = box(s, Inches(0.9), Inches(2.6), Inches(11.6), Inches(4)).text_frame
pairs = [
    ("Selector del panel UMAP",
     "Cuatro botones — GL2Vec / FEATHER / DHC-E / GCN — cambian el espacio del "
     "scatter. Cada subgrafo se reposiciona segun la tecnica elegida (campos "
     "embedGL2Vec / embedFeather / embedDHCE / embedGCN en hotspots.json)."),
    ("Como leerlo",
     "Color por tipo de delito, por numero de delitos o por cluster HDBSCAN. Como el "
     "crimen no forma el layout, ver dos zonas parecidas con criminalidad distinta es "
     "un hallazgo real."),
    ("Reproducir",
     "python embed_subgraphs.py  (recalcula los 3 embeddings sin re-descargar la red vial)."),
]
for i, (k, v) in enumerate(pairs):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    setp(p, k, 19, ACCENT, bold=True, space=2)
    add_para(tf, v, 14, RGBColor(0xC9, 0xCE, 0xD6), space=14)

ref = box(s, Inches(0.9), Inches(6.45), Inches(11.6), Inches(0.9)).text_frame
setp(ref.paragraphs[0],
     "Refs: GL2Vec (ICONIP 2019) · FEATHER (CIKM 2020) · DHC-E (arXiv:2108.02113) · "
     "GCN/ST-GCN (Kipf & Welling 2017; Fan, Hu & Hu, Big Data Cogn. Comput. 2025, 9(7), 179).",
     11, MUTED)

import sys
out = sys.argv[1] if len(sys.argv) > 1 else "Subgrafos_Embeddings_4Tecnicas.pptx"
prs.save(out)
print("OK ->", out)
