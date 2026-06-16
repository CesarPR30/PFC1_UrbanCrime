#!/usr/bin/env python3
"""Genera una presentacion (.pptx) que explica el pipeline de subgrafos:
   generacion, similaridad y embeddings. Basada en preprocess_hotspots.py."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ─── Paleta ──────────────────────────────────────────────────────────────────
INK      = RGBColor(0x1A, 0x1F, 0x2B)   # texto principal / fondo oscuro
ACCENT   = RGBColor(0xFF, 0x6B, 0x4A)   # naranja (acento)
ACCENT2  = RGBColor(0x4A, 0x90, 0xD9)   # azul
MUTED    = RGBColor(0x6B, 0x73, 0x80)   # gris texto
LIGHT    = RGBColor(0xF4, 0xF5, 0xF7)   # gris claro de cajas
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
GREEN    = RGBColor(0x3F, 0xB6, 0x80)

W, H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def box(s, x, y, w, h):
    tb = s.shapes.add_textbox(x, y, w, h)
    tb.text_frame.word_wrap = True
    return tb


def rect(s, x, y, w, h, color, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp


def setp(p, text, size, color, bold=False, align=PP_ALIGN.LEFT, font="Calibri", space=6):
    p.text = text
    p.alignment = align
    p.space_after = Pt(space)
    r = p.runs[0]
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = font
    return p


def add_para(tf, text, size, color, bold=False, align=PP_ALIGN.LEFT, space=6, font="Calibri"):
    p = tf.add_paragraph()
    return setp(p, text, size, color, bold, align, font, space)


def bg(s, color):
    rect(s, 0, 0, W, H, color)


def bullets(tf, items, size=15, color=INK, gap=8):
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        lvl = 0
        txt = it
        if isinstance(it, tuple):
            lvl, txt = it
        p.level = lvl
        setp(p, ("•  " if lvl == 0 else "–  ") + txt, size, color, space=gap)


def header(s, kicker, title):
    rect(s, Inches(0), Inches(0), Inches(0.18), H, ACCENT)
    t = box(s, Inches(0.6), Inches(0.35), Inches(12), Inches(0.4))
    setp(t.text_frame.paragraphs[0], kicker, 13, ACCENT, bold=True)
    t2 = box(s, Inches(0.55), Inches(0.7), Inches(12.2), Inches(0.9))
    setp(t2.text_frame.paragraphs[0], title, 30, INK, bold=True)


# ─── 1. Portada ──────────────────────────────────────────────────────────────
s = slide()
bg(s, INK)
rect(s, Inches(0.9), Inches(2.5), Inches(2.4), Inches(0.09), ACCENT)
t = box(s, Inches(0.85), Inches(2.75), Inches(11.5), Inches(2))
setp(t.text_frame.paragraphs[0], "Subgrafos de criminalidad en Chicago", 40, WHITE, bold=True)
add_para(t.text_frame, "Generacion · Similaridad · Embeddings 2-D", 22, ACCENT, space=4)
t2 = box(s, Inches(0.9), Inches(5.0), Inches(11), Inches(1.5))
setp(t2.text_frame.paragraphs[0],
     "Como se construyen los subgrafos de concentracion criminal sobre la red vial,",
     15, MUTED)
add_para(t2.text_frame, "como se identifican los subgrafos parecidos y como se generan los embeddings.", 15, MUTED)
add_para(t2.text_frame, "preprocess_hotspots.py  ·  embed_subgraphs.py", 13, ACCENT2, space=2)

# ─── 2. Vision general del pipeline ──────────────────────────────────────────
s = slide()
bg(s, WHITE)
header(s, "VISION GENERAL", "El pipeline en 3 etapas")

steps = [
    ("1", "GENERAR", "Clustering greedy BFS\nsobre la red vial", ACCENT,
     "Crimen → nodo → cuenta f(v) → cuencas (basins) → top-20/mes"),
    ("2", "PARECIDOS", "Similaridad estructural\n(coseno)", ACCENT2,
     "Vector de 14 dims (conectividad + geometria) → top-5 'similarTo'"),
    ("3", "EMBEDDINGS", "Topologia + tipo de via\n→ UMAP 2-D", GREEN,
     "Solo forma de la red + jerarquia vial (avenida/calle) → [x, y]"),
]
x = Inches(0.7)
for num, tag, title, col, desc in steps:
    rect(s, x, Inches(2.0), Inches(3.85), Inches(3.6), LIGHT)
    rect(s, x, Inches(2.0), Inches(3.85), Inches(0.12), col)
    c = box(s, x + Inches(0.25), Inches(2.25), Inches(1), Inches(0.8))
    setp(c.text_frame.paragraphs[0], num, 44, col, bold=True)
    tg = box(s, x + Inches(0.25), Inches(3.2), Inches(3.3), Inches(0.4))
    setp(tg.text_frame.paragraphs[0], tag, 16, col, bold=True)
    tt = box(s, x + Inches(0.25), Inches(3.65), Inches(3.4), Inches(1))
    for i, ln in enumerate(title.split("\n")):
        p = tt.text_frame.paragraphs[0] if i == 0 else tt.text_frame.add_paragraph()
        setp(p, ln, 17, INK, bold=True, space=2)
    dd = box(s, x + Inches(0.25), Inches(4.75), Inches(3.4), Inches(0.8))
    setp(dd.text_frame.paragraphs[0], desc, 12.5, MUTED)
    x += Inches(4.1)

f = box(s, Inches(0.7), Inches(6.0), Inches(12), Inches(0.8))
setp(f.text_frame.paragraphs[0],
     "Entrada: crimes.csv + red vial de Chicago (OSMnx).   Salida: public/hotspots.json (top-20 por mes, con similares y embeddings).",
     13, INK)

# ─── 3. Etapa 1 — del crimen al nodo ─────────────────────────────────────────
s = slide()
bg(s, WHITE)
header(s, "ETAPA 1 · GENERACION", "Del crimen al nodo de la red vial")
tf = box(s, Inches(0.7), Inches(1.9), Inches(6.0), Inches(5)).text_frame
bullets(tf, [
    "Se descarga la red vial de Chicago con OSMnx (drive, simplificada).",
    "Cada crimen se ancla (snap) a la arista mas cercana y luego al extremo (nodo) mas proximo.",
    (1, "nearest_edges → se compara distancia haversine a u y a v → gana el mas cercano."),
    "Por mes se cuenta cuantos crimenes caen en cada nodo: f(v).",
    "f(v) es la materia prima del clustering: nodos 'densos' = candidatos a semilla.",
    "Tambien se guarda el conteo por nodo y por mes (node_month_counts) para la historia.",
], size=15)

# diagrama simple
rect(s, Inches(7.1), Inches(2.1), Inches(5.4), Inches(4.4), LIGHT)
d = box(s, Inches(7.35), Inches(2.3), Inches(5), Inches(4)).text_frame
setp(d.paragraphs[0], "Flujo de anclaje", 15, ACCENT, bold=True, space=10)
for ln in ["crimen (lat, lng)", "↓", "arista vial mas cercana", "↓",
           "extremo (nodo) mas cercano", "↓", "f(v) = nº de crimenes en el nodo v",
           "↓", "node_month_counts[v][mes]"]:
    col = INK if "↓" not in ln else MUTED
    bold = "f(v)" in ln or "node_month" in ln
    add_para(d, ln, 15 if "↓" not in ln else 12, col, bold=bold,
             align=PP_ALIGN.CENTER, space=4)

# ─── 4. Etapa 1 — clustering greedy BFS ──────────────────────────────────────
s = slide()
bg(s, WHITE)
header(s, "ETAPA 1 · GENERACION", "Clustering greedy BFS (compute_hotspots)")
tf = box(s, Inches(0.7), Inches(1.85), Inches(12), Inches(5)).text_frame
bullets(tf, [
    "Semillas: se ordenan los nodos por f(v) descendente y se recorren como semillas.",
    "Expansion BFS: desde cada semilla se crece por el grafo vial absorbiendo nodos vecinos con crimen (f ≥ 1), hasta max_hops = 3 saltos.",
    "Puente (bridge = 1): se permite cruzar 1 nodo-conector SIN crimen para alcanzar el siguiente nodo con crimen; ese conector se incorpora a la cuenca.",
    (1, "Sin el puente, una esquina vacia partia una concentracion real en varias cuencas diminutas."),
    "No solapamiento: solo se 'reclaman' los nodos CON crimen (los conectores nunca se reclaman), asi un nucleo ya tomado bloquea las cuencas vecinas.",
    "Cuenca (basin) = nodos con crimen ∪ conectores → subgrafo vial conexo.",
    "Ranking: se ordenan las cuencas por total de crimenes (hyper-volumen) y se conservan las top-20 del mes.",
], size=15, gap=9)

note = box(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.7))
setp(note.text_frame.paragraphs[0],
     "Parametros clave:  top_k = 20   ·   max_hops = 3   ·   bridge = 1",
     14, ACCENT, bold=True)

# ─── 5. Etapa 1 — anatomia de una cuenca ─────────────────────────────────────
s = slide()
bg(s, WHITE)
header(s, "ETAPA 1 · GENERACION", "Anatomia de una cuenca (lo que se guarda)")
items = [
    ("rank", "posicion 1-20 dentro del mes"),
    ("crimes", "total de crimenes en la cuenca"),
    ("peakCrimes", "crimenes en el nodo semilla (el pico)"),
    ("center", "[lat, lng] del nodo semilla"),
    ("nodes", "lista de [lat, lng] de todos los nodos"),
    ("edges", "aristas internas como pares de extremos"),
    ("crimeTypes", "desglose por tipo de delito (para el PieChart)"),
    ("_basinIds", "IDs de nodos — temporal, lo consume el embedding"),
]
y = Inches(1.95)
for k, v in items:
    rect(s, Inches(0.7), y, Inches(3.0), Inches(0.55), INK)
    kb = box(s, Inches(0.8), y + Inches(0.07), Inches(2.8), Inches(0.45))
    setp(kb.text_frame.paragraphs[0], k, 15, WHITE, bold=True, font="Consolas")
    vb = box(s, Inches(3.95), y + Inches(0.05), Inches(8.5), Inches(0.5))
    setp(vb.text_frame.paragraphs[0], v, 14, INK)
    y += Inches(0.6)

# ─── 6. Etapa 2 — subgrafos parecidos ────────────────────────────────────────
s = slide()
bg(s, WHITE)
header(s, "ETAPA 2 · PARECIDOS", "Similaridad estructural (compute_all_similarities)")
tf = box(s, Inches(0.7), Inches(1.85), Inches(6.1), Inches(5)).text_frame
bullets(tf, [
    "Cada cuenca se describe con un vector de 14 dimensiones.",
    "Se normaliza cada columna a [0, 1] (min-max) entre TODAS las cuencas de todos los meses.",
    "Se normaliza cada fila (vector unitario) y se calcula la similitud del coseno entre todos los pares.",
    "Para cada cuenca se guardan sus top-5 vecinos en 'similarTo' (mes, rank, similitud).",
    "Mezclar conectividad + geometria evita que una cadena recta y una en forma de L se traten como gemelas.",
], size=15, gap=9)

rect(s, Inches(7.05), Inches(1.95), Inches(5.5), Inches(4.7), LIGHT)
d = box(s, Inches(7.3), Inches(2.1), Inches(5.1), Inches(4.5)).text_frame
setp(d.paragraphs[0], "Vector de 14 dimensiones", 16, ACCENT2, bold=True, space=8)
setp(add_para(d, "Conectividad (10):", 14, INK, bold=True, space=3), "Conectividad (10):", 14, INK, bold=True, space=3)
for ln in ["log nº nodos · log nº aristas", "densidad · grado medio",
           "pico/total · log intensidad/nodo", "histograma de grados (1/2/3/4+)"]:
    add_para(d, "   " + ln, 13, MUTED, space=2)
add_para(d, "Geometria (4):", 14, INK, bold=True, space=3)
for ln in ["meanEdge (longitud media de arista)", "cvEdge (variabilidad de longitudes)",
           "extent (radio de giro)", "elongation (lineal ↔ blob / forma de L)"]:
    add_para(d, "   " + ln, 13, MUTED, space=2)

# ─── 7. Etapa 3 — embeddings: 2 bloques ──────────────────────────────────────
s = slide()
bg(s, WHITE)
header(s, "ETAPA 3 · EMBEDDINGS", "Solo dos bloques: forma + tipo de via")
blocks = [
    ("TOPOLOGIA (forma)", ACCENT, "peso 1.0",
     ["Firma heat-kernel NetLSD", "(traza del Laplaciano", "normalizado, 16 escalas)",
      "+ log-size, densidad y", "grado medio", "+ geometria metrica", "(distancias y forma)"]),
    ("TIPO DE VIA", ACCENT2, "peso 1.0",
     ["Jerarquia vial de las aristas:", "autopista / avenida /", "colectora / calle / servicio.",
      "Distribucion normalizada", "+ clase dominante.", "Es el dato que separa una",
      "avenida de una calle."]),
]
x = Inches(1.6)
for tag, col, w, lines in blocks:
    rect(s, x, Inches(2.0), Inches(4.6), Inches(4.2), LIGHT)
    rect(s, x, Inches(2.0), Inches(4.6), Inches(0.7), col)
    tg = box(s, x + Inches(0.2), Inches(2.12), Inches(4.2), Inches(0.5))
    setp(tg.text_frame.paragraphs[0], tag, 17, WHITE, bold=True, align=PP_ALIGN.CENTER)
    wb = box(s, x + Inches(0.2), Inches(2.85), Inches(4.2), Inches(0.4))
    setp(wb.text_frame.paragraphs[0], w, 13, col, bold=True, align=PP_ALIGN.CENTER)
    bb = box(s, x + Inches(0.3), Inches(3.4), Inches(4.0), Inches(2.7)).text_frame
    for i, ln in enumerate(lines):
        p = bb.paragraphs[0] if i == 0 else bb.add_paragraph()
        setp(p, ln, 14, INK, space=4)
    x += Inches(4.95)

f = box(s, Inches(0.7), Inches(6.45), Inches(12), Inches(0.9))
setp(f.text_frame.paragraphs[0],
     "El crimen, los POIs y la historia YA NO entran al embedding (siguen guardados para los graficos laterales): "
     "asi el crimen no forma los grupos y queda como hallazgo.",
     13, MUTED)

# ─── 8. Etapa 3 — fusion + UMAP ──────────────────────────────────────────────
s = slide()
bg(s, WHITE)
header(s, "ETAPA 3 · EMBEDDINGS", "Fusion y proyeccion a 2-D con UMAP")
# pipeline horizontal
chain = [("2 bloques\nestandarizados", LIGHT, INK),
         ("× peso\n+ concatenar", LIGHT, INK),
         ("UMAP\n(metric=cosine)", ACCENT, WHITE),
         ("[x, y]\nescalado a [0,1]", GREEN, WHITE)]
x = Inches(0.8)
for i, (txt, col, fg) in enumerate(chain):
    rect(s, x, Inches(2.0), Inches(2.5), Inches(1.3), col)
    b = box(s, x + Inches(0.15), Inches(2.25), Inches(2.2), Inches(0.9))
    for j, ln in enumerate(txt.split("\n")):
        p = b.text_frame.paragraphs[0] if j == 0 else b.text_frame.add_paragraph()
        setp(p, ln, 14, fg, bold=True, align=PP_ALIGN.CENTER, space=2)
    if i < 3:
        ar = box(s, x + Inches(2.5), Inches(2.3), Inches(0.55), Inches(0.7))
        setp(ar.text_frame.paragraphs[0], "→", 30, MUTED, align=PP_ALIGN.CENTER)
    x += Inches(3.05)

tf = box(s, Inches(0.7), Inches(3.7), Inches(12), Inches(3)).text_frame
bullets(tf, [
    "UMAP (n_neighbors=15, min_dist=0.1, metric=cosine, random_state=42) preserva la vecindad local: cuencas parecidas quedan cerca en el scatter.",
    "Si umap-learn no esta disponible, hay un fallback automatico a PCA(2) — el pipeline siempre emite coordenadas.",
    "Se generan 2 espacios (ambos libres de crimen):",
    (1, "topo  → solo la forma de la red (NetLSD + geometria). → embedTopo"),
    (1, "topoRoad  → forma + tipo de via (avenida/calle…). Es el espacio principal/por defecto. → embed2d"),
    "Misma forma pero hecha de avenidas vs. calles residenciales ahora cae en regiones distintas del UMAP.",
], size=14, gap=8)

# ─── 9. Re-embed ligero ──────────────────────────────────────────────────────
s = slide()
bg(s, WHITE)
header(s, "PRACTICO", "Re-calcular embeddings sin re-descargar la red")
tf = box(s, Inches(0.7), Inches(1.95), Inches(12), Inches(5)).text_frame
bullets(tf, [
    "embed_subgraphs.py reutiliza la MISMA logica (compute_subgraph_embeddings) sin volver a bajar la red vial.",
    "Reconstruye la adyacencia desde las aristas ya guardadas en hotspots.json (los IDs de nodo pasan a ser tuplas (lat, lng) redondeadas).",
    "Reconstruye los conteos por nodo y mes anclando cada crimen al nodo de cuenca mas cercano con un KD-tree (radio ≈ 200 m).",
    "Vuelve a calcular embeddings y la similaridad 'similarTo', y reescribe hotspots.json.",
    "Una sola fuente de verdad para los embeddings: el pipeline completo y el re-embed ligero comparten el codigo.",
], size=15, gap=10)
note = box(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.6))
setp(note.text_frame.paragraphs[0], "$ python embed_subgraphs.py", 16, GREEN, bold=True, font="Consolas")

# ─── 10. Cierre ──────────────────────────────────────────────────────────────
s = slide()
bg(s, INK)
rect(s, Inches(0.9), Inches(1.5), Inches(2.4), Inches(0.09), ACCENT)
t = box(s, Inches(0.85), Inches(1.75), Inches(11.5), Inches(1))
setp(t.text_frame.paragraphs[0], "Resumen", 34, WHITE, bold=True)
tf = box(s, Inches(0.9), Inches(2.8), Inches(11.5), Inches(4)).text_frame
pairs = [
    ("Generar", "Greedy BFS sobre la red vial, con puentes para no fragmentar concentraciones reales. Top-20 cuencas/mes por total de crimenes."),
    ("Parecidos", "Vector de 14-dim (conectividad + geometria) + similitud del coseno → top-5 vecinos 'similarTo'."),
    ("Embeddings", "Solo 2 bloques — topologia (NetLSD + geometria) + tipo de via (avenida/calle…) — estandarizados → UMAP 2-D. El crimen no entra: queda como hallazgo."),
]
for i, (k, v) in enumerate(pairs):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    setp(p, k, 20, ACCENT, bold=True, space=2)
    add_para(tf, v, 15, RGBColor(0xC9, 0xCE, 0xD6), space=16)

import sys
out = sys.argv[1] if len(sys.argv) > 1 else "Subgrafos_Generacion_Similaridad_Embeddings.pptx"
prs.save(out)
print("OK ->", out)
